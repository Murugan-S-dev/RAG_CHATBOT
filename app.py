import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
import json
import pytesseract
from PIL import Image
import pandas as pd
from docx import Document
from pptx import Presentation
import mammoth
import datetime

# MY CONTACTS
# GitHub   : [Murugan-S-dev](https://github.com/Murugan-S-dev)
# LinkedIn : [Murugan S](https://www.linkedin.com/in/murugan-s-dev2025)


def get_file_text(uploaded_files):
    text = ""
    for uploaded_file in uploaded_files:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        try:
            if file_extension == 'pdf':
                pdf_reader = PdfReader(uploaded_file)
                for page in pdf_reader.pages:
                    text += page.extract_text() or ""
            elif file_extension in ['png', 'jpg', 'jpeg']:
                image = Image.open(uploaded_file)
                ocr_text = pytesseract.image_to_string(image)
                if not ocr_text.strip():
                    st.warning(
                        f"No text detected in image: {uploaded_file.name}")
                text += ocr_text
            elif file_extension in ['xls', 'xlsx']:
                df = pd.read_excel(uploaded_file)
                text += df.to_string()
            elif file_extension == 'docx':
                doc = Document(uploaded_file)
                for para in doc.paragraphs:
                    text += para.text + '\n'
            elif file_extension == 'doc':
                result = mammoth.convert_to_text(uploaded_file)
                text += result.value
            elif file_extension == 'csv':
                df = pd.read_csv(uploaded_file)
                text += df.to_string()
            elif file_extension == 'txt':
                text += uploaded_file.getvalue().decode("utf-8")
            elif file_extension in ['ppt', 'pptx']:
                prs = Presentation(uploaded_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
        except Exception as e:
            st.warning(f"Error processing {uploaded_file.name}: {e}")
    return text


def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    return text_splitter.split_text(text)


def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    memory = ConversationBufferMemory(
        memory_key='chat_history',
        return_messages=True
    )
    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )


def summarize_text(text):
    summary_prompt = f"Summarize the following document in detail:\n{text[:4000]}"
    llm = ChatOpenAI()
    return llm.predict(summary_prompt)


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']
    latest_response = ""

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            latest_response = message.content
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)

    messages = [{"role": "user" if i % 2 == 0 else "bot", "content": message.content}
                for i, message in enumerate(st.session_state.chat_history)]

    export_data = {
        "timestamp": datetime.datetime.now().isoformat(),
        "chat": messages
    }
    chat_log_json = json.dumps(export_data, indent=2)
    st.download_button("💾 Download Full Chat History (JSON)",
                       chat_log_json, file_name="chat_history.json")

    if latest_response:
        st.download_button("📝 Download Latest Bot Response (TXT)",
                           latest_response, file_name="latest_response.txt")


def main():
    load_dotenv()
    st.set_page_config(page_title="RAG CHATBOT", page_icon="📚")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None
    if "chat_log" not in st.session_state:
        st.session_state.chat_log = []

    st.header("📚 RAG CHATBOT")

    user_question = st.chat_input(
        "Chat with your documents—unlock insights instantly!")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("📂 Upload Documents")
        pdf_docs = st.file_uploader(
            "Upload your files and click 'Analyze'",
            type=['pdf', 'png', 'jpg', 'jpeg', 'xls', 'xlsx',
                  'doc', 'docx', 'csv', 'txt', 'ppt', 'pptx'],
            accept_multiple_files=True
        )

        if st.button("Analyze 🔍"):
            if not pdf_docs:
                st.warning("Please upload documents before analyzing!")
                return

            supported_extensions = {'pdf', 'png', 'jpg', 'jpeg', 'xls',
                                    'xlsx', 'doc', 'docx', 'csv', 'txt', 'ppt', 'pptx'}
            unsupported_files = [f.name for f in pdf_docs if f.name.split(
                '.')[-1].lower() not in supported_extensions]
            if unsupported_files:
                st.warning(
                    f"Unsupported file types detected: {', '.join(unsupported_files)}")
                return

            with st.spinner("Reading and analyzing your documents..."):
                raw_text = get_file_text(pdf_docs)

                if not raw_text.strip():
                    st.error("No readable text found in the uploaded documents.")
                    return

                text_chunks = get_text_chunks(raw_text)
                vectorstore = get_vectorstore(text_chunks)

                st.session_state.chat_history = []
                st.session_state.conversation = get_conversation_chain(
                    vectorstore)

                summary = summarize_text(raw_text)
                st.subheader("📌 Document Summary")
                st.markdown(bot_template.replace(
                    "{{MSG}}", summary), unsafe_allow_html=True)

        st.markdown("---")
        st.subheader("🔁 Reload Chat History")
        uploaded_history = st.file_uploader(
            "Upload previous chat_history.json")
        if uploaded_history:
            data = json.load(uploaded_history)
            st.session_state.chat_log = data.get('chat', [])
            st.success("Chat history reloaded!")

        st.markdown("---")
        if st.button("🗑️ Clear Chat History"):
            st.session_state.conversation = None
            st.session_state.chat_history = None
            st.session_state.chat_log = []
            st.success("Chat history cleared!")

        # 📎 Developer Links

        st.markdown("---")
        st.subheader("👤 Developer Info")

        col1, col2 = st.columns(2)
        with col1:
            if st.button("🔗 GitHub"):
                st.markdown(
                    '<meta http-equiv="refresh" content="0;url=https://github.com/Murugan-S-dev">', unsafe_allow_html=True)
        with col2:
            if st.button("🔗 LinkedIn"):
                st.markdown(
                    '<meta http-equiv="refresh" content="0;url=https://www.linkedin.com/in/murugan-s-dev2025">', unsafe_allow_html=True)

    if st.session_state.chat_log:
        with st.expander("💬 Show Previous Conversation"):
            for message in st.session_state.chat_log:
                if message['role'] == 'user':
                    st.write(user_template.replace(
                        "{{MSG}}", message['content']), unsafe_allow_html=True)
                else:
                    st.write(bot_template.replace(
                        "{{MSG}}", message['content']), unsafe_allow_html=True)


if __name__ == '__main__':
    main()
